"""
반도체 공정 수업자료 문서 처리 시스템
PDF, PPT, DOCX 등의 수업자료를 파싱하여 RAG DB 구축
"""

import os
import json
from typing import List, Dict, Optional
from pathlib import Path
import PyPDF2
from pptx import Presentation
from docx import Document
from openai import AzureOpenAI
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
from azure.search.documents.indexes.models import (
    SearchIndex,
    SimpleField,
    SearchableField,
    SearchField,
    SearchFieldDataType,
    VectorSearch,
    VectorSearchProfile,
    HnswAlgorithmConfiguration,
)
from azure.core.credentials import AzureKeyCredential


class SemiconductorDocumentProcessor:
    """반도체 공정 문서 처리 및 지식 추출"""
    
    def __init__(self):
        # Azure OpenAI 설정
        self.openai_client = AzureOpenAI(
            api_key=os.getenv("AZURE_OPENAI_KEY"),
            api_version="2024-02-15-preview",
            azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT")
        )
        self.gpt_deployment = os.getenv("GPT_DEPLOYMENT_NAME", "gpt-4")
        
        # Azure AI Search 설정
        self.search_endpoint = os.getenv("AZURE_SEARCH_ENDPOINT")
        self.search_key = os.getenv("AZURE_SEARCH_KEY")
        self.index_name = "semiconductor-knowledge"
        
        # 반도체 공정 카테고리
        self.process_categories = [
            "증착 (Deposition)",
            "식각 (Etching)",
            "리소그래피 (Lithography)",
            "이온주입 (Ion Implantation)",
            "확산 (Diffusion)",
            "CMP (Chemical Mechanical Polishing)",
            "세정 (Cleaning)",
            "산화 (Oxidation)",
            "금속화 (Metallization)",
            "패키징 (Packaging)",
            "검사 (Inspection)",
            "이론 (Theory)"
        ]
    
    def parse_pdf(self, pdf_path: str) -> List[Dict]:
        """PDF 파일에서 텍스트 추출 및 구조화"""
        chunks = []
        
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    
                    if text.strip():
                        chunks.append({
                            'content': text,
                            'source': os.path.basename(pdf_path),
                            'page': page_num + 1,
                            'type': 'pdf'
                        })
        except Exception as e:
            print(f"PDF 파싱 오류: {e}")
        
        return chunks
    
    def parse_pptx(self, pptx_path: str) -> List[Dict]:
        """PowerPoint 파일에서 텍스트 추출"""
        chunks = []
        
        try:
            prs = Presentation(pptx_path)
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                # 제목
                if slide.shapes.title:
                    slide_text.append(f"제목: {slide.shapes.title.text}")
                
                # 본문
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    chunks.append({
                        'content': '\n'.join(slide_text),
                        'source': os.path.basename(pptx_path),
                        'page': slide_num + 1,
                        'type': 'pptx'
                    })
        except Exception as e:
            print(f"PPTX 파싱 오류: {e}")
        
        return chunks
    
    def parse_docx(self, docx_path: str) -> List[Dict]:
        """Word 문서에서 텍스트 추출"""
        chunks = []
        
        try:
            doc = Document(docx_path)
            
            current_chunk = []
            for para_num, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    current_chunk.append(para.text)
                    
                    # 500단어마다 청크 분리
                    if len(' '.join(current_chunk).split()) > 500:
                        chunks.append({
                            'content': '\n'.join(current_chunk),
                            'source': os.path.basename(docx_path),
                            'page': para_num // 10 + 1,
                            'type': 'docx'
                        })
                        current_chunk = []
            
            # 남은 텍스트
            if current_chunk:
                chunks.append({
                    'content': '\n'.join(current_chunk),
                    'source': os.path.basename(docx_path),
                    'page': len(chunks) + 1,
                    'type': 'docx'
                })
        except Exception as e:
            print(f"DOCX 파싱 오류: {e}")
        
        return chunks
    
    def extract_semiconductor_knowledge(self, chunks: List[Dict]) -> List[Dict]:
        """
        반도체 공정 지식 추출 및 구조화
        GPT-4로 각 청크를 분석하여 핵심 개념, 공정, 이론 추출
        """
        knowledge_items = []
        
        for chunk in chunks:
            prompt = f"""
다음 반도체 공정 수업자료에서 핵심 지식을 추출하세요:

**원문:**
{chunk['content'][:2000]}

**추출할 정보:**
1. 주요 공정 카테고리 (증착, 식각, 리소그래피 등)
2. 핵심 개념 및 용어
3. 이론적 배경
4. 실무 응용 사례
5. 중요 수식 또는 파라미터
6. 학습 포인트

JSON 형식으로 반환:
{{
    "process_category": "해당 공정 (증착, 식각 등)",
    "key_concepts": ["개념1", "개념2"],
    "theory": "이론적 설명 (2-3문장)",
    "equations": ["수식1", "수식2"] (있는 경우),
    "parameters": ["파라미터1", "파라미터2"],
    "applications": "실무 응용",
    "learning_points": ["학습포인트1", "학습포인트2"],
    "difficulty": "기초/중급/고급"
}}

공정과 관련 없는 내용이면 null 반환.
"""
            
            try:
                response = self.openai_client.chat.completions.create(
                    model=self.gpt_deployment,
                    messages=[
                        {"role": "system", "content": "당신은 반도체 공정 전문가입니다. 수업자료에서 핵심 지식을 추출합니다."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.3,
                    response_format={"type": "json_object"}
                )
                
                knowledge = json.loads(response.choices[0].message.content)
                
                if knowledge.get('process_category'):
                    knowledge['original_content'] = chunk['content']
                    knowledge['source'] = chunk['source']
                    knowledge['page'] = chunk['page']
                    knowledge_items.append(knowledge)
            
            except Exception as e:
                print(f"지식 추출 오류: {e}")
                continue
        
        return knowledge_items
    
    def generate_study_questions(self, knowledge: Dict) -> List[Dict]:
        """
        추출된 지식을 기반으로 학습 질문 자동 생성
        """
        prompt = f"""
다음 반도체 공정 지식을 기반으로 학부생용 학습 질문 5개를 생성하세요:

**공정**: {knowledge.get('process_category', 'N/A')}
**핵심 개념**: {', '.join(knowledge.get('key_concepts', []))}
**이론**: {knowledge.get('theory', 'N/A')}

**질문 유형별로 생성:**
1. 개념 이해 질문 (1개)
2. 원리 설명 질문 (1개)
3. 응용 문제 (1개)
4. 비교/대조 질문 (1개)
5. 실무 상황 질문 (1개)

JSON 배열로 반환:
[
  {{
    "question": "질문 내용",
    "question_type": "개념이해/원리설명/응용/비교/실무",
    "difficulty": "기초/중급/고급",
    "answer": "모범 답변 (2-3문장)",
    "keywords": ["키워드1", "키워드2"],
    "related_concepts": ["관련개념1", "관련개념2"]
  }}
]
"""
        
        try:
            response = self.openai_client.chat.completions.create(
                model=self.gpt_deployment,
                messages=[
                    {"role": "system", "content": "당신은 반도체 공학 교수입니다. 효과적인 학습 질문을 만듭니다."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                response_format={"type": "json_object"}
            )
            
            result = json.loads(response.choices[0].message.content)
            questions = result if isinstance(result, list) else result.get('questions', [])
            
            # 메타데이터 추가
            for q in questions:
                q['process_category'] = knowledge.get('process_category')
                q['source'] = knowledge.get('source')
            
            return questions
        
        except Exception as e:
            print(f"질문 생성 오류: {e}")
            return []
    
    def get_embedding(self, text: str) -> List[float]:
        """텍스트 임베딩 생성"""
        response = self.openai_client.embeddings.create(
            model="text-embedding-ada-002",
            input=text
        )
        return response.data[0].embedding
    
    def create_search_index(self):
        """반도체 지식 검색 인덱스 생성"""
        index_client = SearchIndexClient(
            endpoint=self.search_endpoint,
            credential=AzureKeyCredential(self.search_key)
        )
        
        # 벡터 검색 설정
        vector_search = VectorSearch(
            algorithms=[
                HnswAlgorithmConfiguration(
                    name="hnsw-config",
                    parameters={
                        "m": 4,
                        "efConstruction": 400,
                        "efSearch": 500,
                        "metric": "cosine"
                    }
                )
            ],
            profiles=[
                VectorSearchProfile(
                    name="vector-profile",
                    algorithm_configuration_name="hnsw-config",
                )
            ]
        )
        
        # 필드 정의
        fields = [
            SimpleField(name="id", type=SearchFieldDataType.String, key=True),
            SearchableField(name="question", type=SearchFieldDataType.String),
            SearchableField(name="answer", type=SearchFieldDataType.String),
            SearchableField(name="process_category", type=SearchFieldDataType.String, filterable=True),
            SearchableField(name="question_type", type=SearchFieldDataType.String, filterable=True),
            SimpleField(name="difficulty", type=SearchFieldDataType.String, filterable=True),
            SearchableField(name="theory", type=SearchFieldDataType.String),
            SearchableField(name="source", type=SearchFieldDataType.String),
            SearchField(
                name="contentVector",
                type=SearchFieldDataType.Collection(SearchFieldDataType.Single),
                searchable=True,
                vector_search_dimensions=1536,
                vector_search_profile_name="vector-profile"
            ),
            SearchableField(
                name="keywords",
                type=SearchFieldDataType.Collection(SearchFieldDataType.String),
                filterable=True
            ),
            SearchableField(
                name="related_concepts",
                type=SearchFieldDataType.Collection(SearchFieldDataType.String)
            )
        ]
        
        # 인덱스 생성
        index = SearchIndex(
            name=self.index_name,
            fields=fields,
            vector_search=vector_search
        )
        
        result = index_client.create_or_update_index(index)
        print(f"인덱스 '{result.name}' 생성 완료!")
        return result
    
    def upload_to_search(self, questions: List[Dict]) -> Dict:
        """생성된 질문을 검색 인덱스에 업로드"""
        search_client = SearchClient(
            endpoint=self.search_endpoint,
            index_name=self.index_name,
            credential=AzureKeyCredential(self.search_key)
        )
        
        # 시작 ID 확인
        try:
            results = search_client.search(search_text="*", top=1, select=["id"])
            existing_ids = [int(doc['id']) for doc in results if doc['id'].isdigit()]
            start_id = max(existing_ids) + 1 if existing_ids else 1
        except:
            start_id = 1
        
        # 문서 준비
        documents = []
        for i, q in enumerate(questions):
            # 임베딩 생성
            embedding_text = f"{q['question']} {q.get('answer', '')} {' '.join(q.get('keywords', []))}"
            embedding = self.get_embedding(embedding_text)
            
            doc = {
                "id": str(start_id + i),
                "question": q['question'],
                "answer": q.get('answer', ''),
                "process_category": q.get('process_category', '이론'),
                "question_type": q.get('question_type', '개념이해'),
                "difficulty": q.get('difficulty', '중급'),
                "theory": q.get('theory', ''),
                "source": q.get('source', ''),
                "contentVector": embedding,
                "keywords": q.get('keywords', []),
                "related_concepts": q.get('related_concepts', [])
            }
            documents.append(doc)
        
        # 업로드
        result = search_client.upload_documents(documents=documents)
        
        success_count = sum(1 for r in result if r.succeeded)
        return {
            'success': success_count,
            'failed': len(result) - success_count,
            'total': len(documents)
        }
    
    def process_course_materials(self, file_paths: List[str]) -> Dict:
        """
        수업자료 일괄 처리
        
        Args:
            file_paths: 수업자료 파일 경로 리스트
        
        Returns:
            처리 결과 통계
        """
        all_chunks = []
        
        for file_path in file_paths:
            ext = Path(file_path).suffix.lower()
            
            if ext == '.pdf':
                chunks = self.parse_pdf(file_path)
            elif ext == '.pptx':
                chunks = self.parse_pptx(file_path)
            elif ext == '.docx':
                chunks = self.parse_docx(file_path)
            else:
                print(f"지원하지 않는 파일 형식: {ext}")
                continue
            
            all_chunks.extend(chunks)
        
        print(f"총 {len(all_chunks)}개 청크 추출")
        
        # 지식 추출
        knowledge_items = self.extract_semiconductor_knowledge(all_chunks)
        print(f"총 {len(knowledge_items)}개 지식 항목 추출")
        
        # 학습 질문 생성
        all_questions = []
        for knowledge in knowledge_items:
            questions = self.generate_study_questions(knowledge)
            all_questions.extend(questions)
        
        print(f"총 {len(all_questions)}개 질문 생성")
        
        # 인덱스 생성 및 업로드
        self.create_search_index()
        upload_result = self.upload_to_search(all_questions)
        
        return {
            'files_processed': len(file_paths),
            'chunks_extracted': len(all_chunks),
            'knowledge_items': len(knowledge_items),
            'questions_generated': len(all_questions),
            'upload_result': upload_result
        }


# 테스트 코드
if __name__ == "__main__":
    processor = SemiconductorDocumentProcessor()
    
    # 샘플 사용
    print("반도체 공정 문서 처리 시스템")
    print("\n사용 예시:")
    print("""
# 수업자료 처리
processor = SemiconductorDocumentProcessor()

# 여러 파일 일괄 처리
file_paths = [
    'lecture01_deposition.pdf',
    'lecture02_etching.pptx',
    'lecture03_lithography.docx'
]

result = processor.process_course_materials(file_paths)
print(f"처리 완료: {result}")
    """)
